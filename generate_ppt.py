import json
import os
import google.generativeai as genai
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def get_api_key():
    try:
        with open('.env', 'r') as f:
            return f.read().strip()
    except FileNotFoundError:
        print("Error: .env file not found.")
        return None

def generate_image(prompt, output_path, api_key):
    try:
        genai.configure(api_key=api_key)
        model = genai.GenerativeModel('gemini-3-pro-image-preview')
        
        print(f"Requesting image for prompt: {prompt[:30]}...")
        response = model.generate_content(prompt)
        
        if response.parts:
            for part in response.parts:
                if hasattr(part, 'inline_data') and part.inline_data:
                    # Found image data
                    with open(output_path, 'wb') as f:
                        f.write(part.inline_data.data)
                    return True
        
        print("No image found in response.")
        return False
            
    except Exception as e:
        print(f"Failed to generate image for prompt: {prompt[:30]}... Error: {e}")
        return False

from datetime import datetime

def create_presentation(json_file='slides.json', output_file_base='nano_banana_presentation'):
    api_key = get_api_key()
    if not api_key:
        print("Skipping image generation due to missing API key.")

    try:
        with open(json_file, 'r', encoding='utf-8') as f:
            data = json.load(f)
            # Handle both formats: direct array or object with 'slides' key
            if isinstance(data, dict) and 'slides' in data:
                slides_data = data['slides']
            elif isinstance(data, list):
                slides_data = data
            else:
                print(f"Error: Unexpected JSON format in {json_file}")
                return
    except FileNotFoundError:
        print(f"Error: {json_file} not found.")
        return

    # Generate a timestamp for the output file
    timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
    output_file = f"{output_file_base}_{timestamp}.pptx"

    prs = Presentation()
    # Set slide dimensions to 16:9 aspect ratio
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Create images directory if it doesn't exist
    if not os.path.exists('generated_images'):
        os.makedirs('generated_images')

    for i, slide_data in enumerate(slides_data):
        # Use a blank layout for custom positioning
        slide_layout = prs.slide_layouts[6] 
        slide = prs.slides.add_slide(slide_layout)

        # Title
        title_left = Inches(0.5)
        title_top = Inches(0.3)
        title_width = Inches(12)
        title_height = Inches(1.0)
        
        title_box = slide.shapes.add_textbox(title_left, title_top, title_width, title_height)
        title_tf = title_box.text_frame
        title_tf.text = slide_data.get('title', 'No Title')
        title_tf.paragraphs[0].font.size = Pt(40)
        title_tf.paragraphs[0].font.bold = True

        # Content (Text) - Left side
        left = Inches(0.5)
        top = Inches(1.5)
        width = Inches(6.0)
        height = Inches(5.0)
        
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        
        # Handle content as either string or list
        content = slide_data.get('content', '')
        if isinstance(content, list):
            # Convert list to bullet points
            content_text = '\n'.join([f"• {item}" for item in content])
        else:
            content_text = content
            
        tf.text = content_text
        tf.word_wrap = True
        
        # Adjust font size for content
        for paragraph in tf.paragraphs:
            paragraph.font.size = Pt(20)

        # Image - Right side
        img_prompt = slide_data.get('image_prompt', '')
        img_filename = f"generated_images/slide_{i+1}.png"
        
        image_generated = False
        if api_key and img_prompt:
            print(f"Generating image for slide {i+1}...")
            image_generated = generate_image(img_prompt, img_filename, api_key)
        
        if image_generated and os.path.exists(img_filename):
            img_left = Inches(7.0)
            img_top = Inches(1.5)
            img_width = Inches(5.8)
            # Add image
            slide.shapes.add_picture(img_filename, img_left, img_top, width=img_width)
        else:
            # Placeholder if no image
            placeholder_left = Inches(7.0)
            placeholder_top = Inches(1.5)
            placeholder_width = Inches(5.8)
            placeholder_height = Inches(4.0)
            
            shape = slide.shapes.add_shape(
                1, # msoShapeRectangle
                placeholder_left, placeholder_top, placeholder_width, placeholder_height
            )
            # Make it look like a placeholder (light gray fill, border)
            shape.fill.solid()
            shape.fill.fore_color.rgb = RGBColor(238, 238, 238) # Light gray (0xEEEEEE)
            shape.line.color.rgb = RGBColor(136, 136, 136) # Darker gray border (0x888888)
            
            # Add text to the placeholder
            p_tf = shape.text_frame
            p_tf.text = f"Image Placeholder\n\nPrompt:\n{img_prompt}" if img_prompt else "No Image Prompt"
            p_tf.paragraphs[0].alignment = PP_ALIGN.CENTER
            p_tf.paragraphs[0].font.bold = True
        
        # Add notes
        notes_slide = slide.notes_slide
        text_frame = notes_slide.notes_text_frame
        text_frame.text = f"Image Prompt: {img_prompt}"

    prs.save(output_file)
    print(f"Presentation saved to {output_file}")

if __name__ == "__main__":
    create_presentation()
